from pptx import Presentation
import os


def extract_images_from_pptx(pptx_path, output_dir):
    """
    Extract all images from a PowerPoint file and save them to a directory.
    """
    prs = Presentation(pptx_path)
    image_count = 0

    # Create output directory if it doesn't exist
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    # Go through all slides
    for slide in prs.slides:
        # Go through all shapes in the slide
        for shape in slide.shapes:
            # Check if shape has an image
            if hasattr(shape, "image"):
                # Get image bytes
                image = shape.image
                # Get image file extension
                image_ext = image.ext

                # Create image file path
                image_path = os.path.join(
                    output_dir, f'image_{image_count}{image_ext}')

                # Save image
                with open(image_path, 'wb') as f:
                    f.write(image.blob)

                print(f"Saved image: {image_path}")
                image_count += 1

    print(f"Extracted {image_count} images")

# Example usage:
# extract_images_from_pptx("aws_services.pptx", "output_directory")
